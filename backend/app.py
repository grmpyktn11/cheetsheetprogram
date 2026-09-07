import os
import base64
import uuid
from dotenv import load_dotenv
from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from fpdf import FPDF
import anthropic


# Load environment variables
load_dotenv()

# Initialize the Claude client. The zero-arg constructor resolves credentials
# from the environment: ANTHROPIC_API_KEY, or the profile stored by `ant auth
# login`. Nothing to hardcode and nothing to keep in .env.
client = anthropic.Anthropic()

MODEL = "claude-opus-5"

# Flask app configuration
app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = "./uploads"
OUTPUT_PDF = "cheat_sheet.pdf"   # default for create_pdf() called outside the route
FONT_PATH = "indie.ttf"


def ensure_upload_folder():
    """Create uploads folder if it doesn't exist."""
    if not os.path.exists(UPLOAD_FOLDER):
        os.makedirs(UPLOAD_FOLDER)


# At import, not under __main__. Only `python app.py` ran this before, so under
# a WSGI server - which is how it is actually deployed - the folder was never
# created and the first upload failed on file.save().
ensure_upload_folder()


def extract_text_from_pptx(file_path):
    """
    Extract all text content from a PowerPoint presentation.
    
    Args:
        file_path: Path to the .pptx file
        
    Returns:
        Extracted text with bullet points converted to '>' and newlines removed
    """
    prs = Presentation(file_path)
    text_content = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content.append(shape.text_frame.text)
    
    # Join and clean the text
    full_text = ''.join(text_content)
    full_text = full_text.replace("•", ">").replace('\n', ' ')
    
    return full_text


SUMMARY_SYSTEM_PROMPT = """You read a whole PowerPoint and write a one to two page cheat sheet from it. Format it as plain text that can be written straight into a PDF: use bullet points and clear structure, and do not use Markdown syntax such as ** or ##, since it is not rendered and would be printed literally. Return only the cheat sheet itself, with no preamble."""


def generate_summary(content):
    """
    Generate a concise summary using Claude.
    
    Args:
        content: Text content to summarize
        
    Returns:
        Summarized text formatted for PDF output
    """
    response = client.messages.create(
        model=MODEL,
        max_tokens=16000,
        system=SUMMARY_SYSTEM_PROMPT,
        messages=[{"role": "user", "content": content}],
    )
    
    # response.content is a list of blocks, not a string - a thinking block can
    # precede the text, so filter by type rather than indexing content[0].
    return "".join(
        block.text for block in response.content if block.type == "text"
    )


def create_pdf(content, output_path=OUTPUT_PDF):
    """
    Generate a PDF document from text content.
    
    Args:
        content: Text content to include in the PDF
        output_path: Path where the PDF should be saved
    """
    pdf = FPDF()
    pdf.add_page()
    
    # Add custom font
    pdf.add_font('CustomFont', '', FONT_PATH, uni=True)
    
    # Title
    pdf.set_font('CustomFont', '', 14)
    pdf.cell(0, 10, txt='CHEAT SHEET', align='C', ln=1)
    
    # Content
    pdf.set_font('CustomFont', '', 13)
    pdf.multi_cell(0, 10, txt=content, border=0, align='L', fill=False)
    
    pdf.output(output_path)


def encode_file_to_base64(file_path):
    """
    Convert a file to base64 encoding.
    
    Args:
        file_path: Path to the file to encode
        
    Returns:
        Base64 encoded string
    """
    with open(file_path, 'rb') as file:
        file_bytes = base64.b64encode(file.read())
        return file_bytes.decode('ascii')


@app.route("/upload", methods=["POST"])
def upload_file():
    """
    Handle file upload endpoint.
    Accepts a PowerPoint file, extracts text, generates summary, and returns PDF.
    """
    # Validate file in request
    if "file" not in request.files:
        return jsonify({"error": "No file part in request"}), 400

    file = request.files["file"]

    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400
    
    if not file.filename.endswith('.pptx'):
        return jsonify({"error": "Only .pptx files are supported"}), 400

    # Both paths used to be fixed names - every request wrote presentation.pptx
    # and cheat_sheet.pdf - so two people uploading at once each got back
    # whichever summary finished last. A token per request keeps them apart, and
    # both files are removed once the PDF is encoded into the response.
    token = uuid.uuid4().hex
    file_path = os.path.join(UPLOAD_FOLDER, f"{token}.pptx")
    pdf_path = os.path.join(UPLOAD_FOLDER, f"{token}.pdf")

    try:
        file.save(file_path)

        # Process the file
        text = extract_text_from_pptx(file_path)
        summary = generate_summary(text)
        create_pdf(summary, pdf_path)

        # Encode PDF and return
        encoded_pdf = encode_file_to_base64(pdf_path)

        return jsonify({'pdfEncoded': encoded_pdf}), 200

    except Exception as e:
        return jsonify({"error": f"Processing failed: {str(e)}"}), 500

    finally:
        for path in (file_path, pdf_path):
            try:
                os.remove(path)
            except OSError:
                pass


@app.route("/health", methods=["GET"])
def health_check():
    """Health check endpoint."""
    return jsonify({"status": "healthy"}), 200


if __name__ == "__main__":
    ensure_upload_folder()
    app.run(debug=True)
