from pptx import Presentation
file_path = "uploads/pp.pptx"
prs = Presentation(file_path)

def getText():
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text += shape.text_frame.text
            
    text = text.replace('\n', '')   
    print(text)


